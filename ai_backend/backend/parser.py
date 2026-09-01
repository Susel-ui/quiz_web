from __future__ import annotations

"""Extract text from uploaded learning material before it reaches the LLM.

The quiz generator should reason over plain text, not raw files. This module
keeps extraction concerns isolated so API handlers and Streamlit code can share
one predictable contract for PDFs, presentations, and simple text documents.
"""

from io import BytesIO
import os
import re
import zipfile
from typing import List, Optional

from pydantic import BaseModel, Field


class ExtractionResult(BaseModel):
    filename: Optional[str] = None
    file_type: str
    text: str
    unit_count: int = 0
    per_unit_text: List[str] = Field(default_factory=list)
    warnings: List[str] = Field(default_factory=list)

    @property
    def char_count(self) -> int:
        return len(self.text)

    @property
    def is_empty(self) -> bool:
        return not self.text.strip()


def _read_file_bytes(file_obj) -> bytes:
    if isinstance(file_obj, bytes):
        return file_obj
    if isinstance(file_obj, bytearray):
        return bytes(file_obj)
    if hasattr(file_obj, "read"):
        data = file_obj.read()
        if hasattr(file_obj, "seek"):
            try:
                file_obj.seek(0)
            except Exception:
                pass
        return data
    raise TypeError("file_obj must be bytes or a readable file-like object")


def _extension(filename: Optional[str]) -> str:
    if not filename:
        return ""
    return os.path.splitext(filename)[1].lower().lstrip(".")


def _normalise_text(parts: List[str]) -> str:
    joined = "\n\n".join(part.strip() for part in parts if part and part.strip())
    return re.sub(r"\n{3,}", "\n\n", joined).strip()


def _extract_pdf(data: bytes) -> tuple[str, List[str], List[str]]:
    warnings: List[str] = []
    try:
        import pdfplumber
    except ImportError as exc:
        raise RuntimeError("PDF extraction requires pdfplumber to be installed") from exc

    per_page: List[str] = []
    with pdfplumber.open(BytesIO(data)) as pdf:
        for index, page in enumerate(pdf.pages, start=1):
            page_text = page.extract_text() or ""
            if not page_text.strip():
                warnings.append(f"Page {index} had no extractable text.")
            per_page.append(page_text.strip())
    return _normalise_text(per_page), per_page, warnings


def _extract_pptx(data: bytes) -> tuple[str, List[str], List[str]]:
    warnings: List[str] = []
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("PowerPoint extraction requires python-pptx to be installed") from exc

    presentation = Presentation(BytesIO(data))
    per_slide: List[str] = []
    for slide in presentation.slides:
        chunks: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                chunks.append(shape.text)
        per_slide.append("\n".join(chunks).strip())
    return _normalise_text(per_slide), per_slide, warnings


def _extract_docx(data: bytes) -> tuple[str, List[str], List[str]]:
    warnings: List[str] = []
    try:
        with zipfile.ZipFile(BytesIO(data)) as docx_zip:
            xml = docx_zip.read("word/document.xml").decode("utf-8", errors="ignore")
    except Exception as exc:
        raise RuntimeError("Could not read DOCX text from word/document.xml") from exc

    paragraphs = re.findall(r"<w:p\b[^>]*>(.*?)</w:p>", xml, flags=re.DOTALL)
    texts: List[str] = []
    for paragraph in paragraphs:
        text_runs = re.findall(r"<w:t(?:\s[^>]*)?>(.*?)</w:t>", paragraph, flags=re.DOTALL)
        if text_runs:
            text = "".join(text_runs)
            text = re.sub(r"<[^>]+>", "", text)
            texts.append(text)
    return _normalise_text(texts), texts, warnings


def _decode_text(data: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "cp1252", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="ignore")


def extract_text(file_obj, filename: str = None) -> ExtractionResult:
    data = _read_file_bytes(file_obj)
    file_type = _extension(filename)
    warnings: List[str] = []

    if not data:
        return ExtractionResult(
            filename=filename,
            file_type=file_type or "unknown",
            text="",
            unit_count=0,
            per_unit_text=[],
            warnings=["Uploaded file was empty."],
        )

    if file_type == "pdf":
        text, per_unit_text, warnings = _extract_pdf(data)
    elif file_type in {"pptx", "ppt"}:
        if file_type == "ppt":
            warnings.append("Legacy .ppt files may not extract reliably; prefer .pptx.")
        text, per_unit_text, ppt_warnings = _extract_pptx(data)
        warnings.extend(ppt_warnings)
    elif file_type == "docx":
        text, per_unit_text, warnings = _extract_docx(data)
    else:
        text = _decode_text(data)
        per_unit_text = [text]
        if not file_type:
            file_type = "text"

    return ExtractionResult(
        filename=filename,
        file_type=file_type or "unknown",
        text=text.strip(),
        unit_count=len(per_unit_text),
        per_unit_text=per_unit_text,
        warnings=warnings,
    )


def truncate_for_llm(text: str, max_chars: int = 18000) -> str:
    if max_chars <= 0:
        raise ValueError("max_chars must be greater than zero")
    if len(text) <= max_chars:
        return text
    return text[:max_chars].rsplit(" ", 1)[0].strip()
